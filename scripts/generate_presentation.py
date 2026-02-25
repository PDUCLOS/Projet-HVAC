"""Generate a professional PowerPoint presentation for the HVAC Market Analysis project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- Color palette --
DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
MEDIUM_BLUE = RGBColor(0x1E, 0x88, 0xE5)
LIGHT_BLUE = RGBColor(0xBB, 0xDE, 0xFB)
ACCENT_ORANGE = RGBColor(0xFF, 0x6F, 0x00)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x42, 0x42, 0x42)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=BLACK, spacing=Pt(6)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_kpi_card(slide, left, top, width, height, value, label, color=MEDIUM_BLUE):
    """Add a KPI card with value and label."""
    card = add_shape_rect(slide, left, top, width, height, WHITE)
    card.shadow.inherit = False

    # Colored top bar
    add_shape_rect(slide, left, top, width, Pt(5), color)

    # Value
    add_text_box(slide, left, top + Inches(0.25), width, Inches(0.6),
                 value, font_size=28, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Label
    add_text_box(slide, left, top + Inches(0.85), width, Inches(0.5),
                 label, font_size=12, color=GRAY,
                 alignment=PP_ALIGN.CENTER)


def add_table(slide, left, top, width, height, rows, cols, data, header_color=DARK_BLUE):
    """Add a formatted table."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = data[row_idx][col_idx]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Calibri"
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = BLACK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return table_shape


# =============================================================================
# CREATE PRESENTATION
# =============================================================================

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

blank_layout = prs.slide_layouts[6]  # Blank layout

# =========================================================================
# SLIDE 1 — Title
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BLUE)

# Decorative bar at top
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_ORANGE)

# Main title
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "HVAC Market Analysis", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "Predictive Analysis of the HVAC Market across Metropolitan France",
             font_size=22, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

# Separator line
add_shape_rect(slide, Inches(5), Inches(4.0), Inches(3), Pt(3), ACCENT_ORANGE)

# Author
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             "Patrice DUCLOS  |  Senior Data Analyst  |  20 years of experience",
             font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

# Certification
add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
             "Data Science Lead Certification  —  Jedha Bootcamp (Bac+5, RNCP Level 7)",
             font_size=14, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

# Date
add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
             "February 2026", font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 2 — Agenda
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(8), Inches(0.7),
             "Agenda", font_size=32, color=WHITE, bold=True)

items = [
    ("01", "Project Vision & Objectives"),
    ("02", "Data Sources & Collection"),
    ("03", "Technical Architecture"),
    ("04", "Data Pipeline (ELT)"),
    ("05", "Feature Engineering"),
    ("06", "Machine Learning Models & Results"),
    ("07", "API & Dashboard"),
    ("08", "Deployment & DevOps"),
    ("09", "Certification Coverage"),
    ("10", "Key Takeaways & Next Steps"),
]

for i, (num, title) in enumerate(items):
    y = Inches(1.5) + Inches(i * 0.55)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = MEDIUM_BLUE if i % 2 == 0 else ACCENT_ORANGE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False

    add_text_box(slide, Inches(1.8), y + Inches(0.02), Inches(8), Inches(0.4),
                 title, font_size=18, color=DARK_GRAY)

# =========================================================================
# SLIDE 3 — Project Vision
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "01  Project Vision & Objectives", font_size=32, color=WHITE, bold=True)

# Mission statement
add_shape_rect(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.1), LIGHT_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.5), Inches(11.3), Inches(0.9),
             "Build a complete end-to-end data pipeline to analyze and predict HVAC equipment "
             "installations (heat pumps, air conditioning) across 96 metropolitan departments "
             "by cross-referencing energy, weather, economic, and construction data.",
             font_size=16, color=DARK_BLUE)

# Objectives
add_text_box(slide, Inches(0.8), Inches(2.8), Inches(5), Inches(0.5),
             "Objectives", font_size=22, color=DARK_BLUE, bold=True)

objectives = [
    "Collect & merge 6 open data sources (ADEME, Open-Meteo, INSEE, Eurostat, SITADEL)",
    "Build ~100+ ML features with domain expertise (PAC efficiency, altitude, climate)",
    "Train & compare ML models: Ridge, LightGBM, Prophet, LSTM",
    "Serve predictions via REST API (FastAPI) + interactive dashboard (Streamlit)",
    "Deploy with Docker, Kubernetes manifests, Airflow orchestration",
    "Demonstrate end-to-end Data Science skills for Bac+5 certification",
]
add_bullet_list(slide, Inches(1.0), Inches(3.3), Inches(11), Inches(3.5),
                [f"  {obj}" for obj in objectives], font_size=15, color=DARK_GRAY)

# Design principles box
add_shape_rect(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.9), LIGHT_GRAY)
add_text_box(slide, Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.7),
             "Design Principles:  100% local-first (SQLite)  |  Portable & reproducible  |  "
             "All Open Data (no API key)  |  Professional-grade code & testing",
             font_size=14, color=GRAY)

# =========================================================================
# SLIDE 4 — Data Sources
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "02  Data Sources & Collection", font_size=32, color=WHITE, bold=True)

data_sources = [
    ["Source", "API", "Coverage", "Data", "Volume"],
    ["DPE ADEME", "data.ademe.fr", "96 departments", "Energy diagnostics (PAC detection)", "~1.4M rows"],
    ["Open-Meteo", "archive-api.open-meteo.com", "96 prefectures, 7 years", "Temperature, HDD/CDD, frost", "~250K rows"],
    ["INSEE BDM", "bdm.insee.fr (SDMX)", "France, monthly", "Confidence, IPI, business climate", "~500 rows"],
    ["Eurostat", "eurostat package", "France, monthly", "Industrial Production Index HVAC", "~200 rows"],
    ["SITADEL", "DiDo API (SDES)", "96 depts, monthly", "Building permits", "~5K rows"],
    ["INSEE Filosofi", "Reference CSV", "96 depts (static)", "Income, price/m2, housing stock", "96 rows"],
]

add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.2),
          len(data_sources), 5, data_sources)

# Highlight box
add_shape_rect(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.0), LIGHT_BLUE)
add_text_box(slide, Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.8),
             "All sources are 100% Open Data — no API key required. "
             "DPE ADEME serves as the target variable: each diagnostic mentioning a heat pump "
             "or air conditioning system is counted as a recent installation.",
             font_size=14, color=DARK_BLUE)

# Key insight
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(5), Inches(0.4),
             "Plugin Architecture", font_size=18, color=DARK_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.5),
             "Each collector inherits from BaseCollector with auto-registration. "
             "Adding a new source = 1 file, ~50 lines, zero config change.",
             font_size=14, color=DARK_GRAY)

# =========================================================================
# SLIDE 5 — Architecture Overview
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "03  Technical Architecture", font_size=32, color=WHITE, bold=True)

# Pipeline flow (boxes with arrows)
stages = [
    ("6 Sources\nOpen Data", MEDIUM_BLUE),
    ("Collection\n(5 Collectors)", ACCENT_ORANGE),
    ("Data Lake\nCSV / SQLite", RGBColor(0x7B, 0x1F, 0xA2)),
    ("ETL\nClean + Merge", ACCENT_RED),
    ("Features\n~100 variables", ACCENT_GREEN),
    ("ML Models\nRidge, LightGBM", RGBColor(0xF9, 0xA8, 0x25)),
    ("API + Dashboard\nFastAPI + Streamlit", MEDIUM_BLUE),
]

box_w = Inches(1.55)
box_h = Inches(1.0)
start_x = Inches(0.5)
y_pos = Inches(1.5)
gap = Inches(0.25)

for i, (label, color) in enumerate(stages):
    x = start_x + i * (box_w + gap)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, box_w, box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    for line_idx, line in enumerate(label.split("\n")):
        if line_idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = (line_idx == 0)
        p.alignment = PP_ALIGN.CENTER

    # Arrow
    if i < len(stages) - 1:
        arrow_x = x + box_w
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, arrow_x, y_pos + Inches(0.35),
            gap, Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY
        arrow.line.fill.background()

# Tech stack table
add_text_box(slide, Inches(0.8), Inches(3.0), Inches(5), Inches(0.4),
             "Tech Stack", font_size=20, color=DARK_BLUE, bold=True)

stack_data = [
    ["Category", "Technologies"],
    ["Language", "Python 3.10+"],
    ["Data", "pandas, numpy, SQLAlchemy"],
    ["ML / DL", "scikit-learn, LightGBM, Prophet, PyTorch (LSTM), SHAP"],
    ["API", "FastAPI, Pydantic, uvicorn"],
    ["Dashboard", "Streamlit, Plotly"],
    ["Database", "SQLite, PostgreSQL, SQL Server"],
    ["DevOps", "Docker, Kubernetes, Airflow, Makefile"],
    ["Cloud", "pCloud (sync), Render.com (deploy)"],
]

add_table(slide, Inches(0.6), Inches(3.5), Inches(12.1), Inches(3.5),
          len(stack_data), 2, stack_data)

# =========================================================================
# SLIDE 6 — Data Pipeline (ELT)
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "04  Data Pipeline (ELT)", font_size=32, color=WHITE, bold=True)

# ELT justification
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
             "ELT Pattern: Load raw first, then Transform", font_size=18,
             color=DARK_BLUE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.6),
             "Raw CSVs in data/raw/ are the source of truth. If a cleaning rule changes, "
             "re-run 'clean' without re-collecting. Each step is idempotent and reproducible.",
             font_size=14, color=DARK_GRAY)

# Pipeline steps
pipeline_steps = [
    ("1. Collect", "5 APIs  ->  data/raw/", "15-45 min", MEDIUM_BLUE),
    ("2. Clean", "Dedup, NaN, clipping, type fix", "< 2 min", ACCENT_ORANGE),
    ("3. Merge", "Multi-source JOIN (date x dept)", "< 1 min", RGBColor(0x7B, 0x1F, 0xA2)),
    ("4. Features", "Lags, rolling, interactions", "< 30 sec", ACCENT_GREEN),
    ("5. Outliers", "IQR + Z-score + Isolation Forest", "< 30 sec", ACCENT_RED),
    ("6. Train", "Ridge, LightGBM, Prophet, LSTM", "1-5 min", RGBColor(0xF9, 0xA8, 0x25)),
    ("7. Evaluate", "RMSE, MAE, SHAP, visualizations", "1-3 min", MEDIUM_BLUE),
]

for i, (step, desc, time, color) in enumerate(pipeline_steps):
    y = Inches(2.5) + i * Inches(0.6)
    # Step badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.8), y, Inches(2.0), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = step
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Description
    add_text_box(slide, Inches(3.0), y, Inches(6.5), Inches(0.45),
                 desc, font_size=14, color=DARK_GRAY)

    # Time
    add_text_box(slide, Inches(9.8), y, Inches(2), Inches(0.45),
                 time, font_size=12, color=GRAY, alignment=PP_ALIGN.RIGHT)

# CLI
add_shape_rect(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5), LIGHT_GRAY)
add_text_box(slide, Inches(1.0), Inches(6.85), Inches(11.3), Inches(0.4),
             "CLI:  python -m src.pipeline all  |  Interactive mode:  python -m src.pipeline clean -i",
             font_size=13, color=GRAY)

# =========================================================================
# SLIDE 7 — Feature Engineering
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "05  Feature Engineering (~100+ features)", font_size=32, color=WHITE, bold=True)

# Feature categories
feat_data = [
    ["Category", "Count", "Examples", "Purpose"],
    ["Temporal", "7", "month_sin/cos, is_heating, is_cooling", "Seasonal patterns"],
    ["Lags (1/3/6 mo)", "27", "nb_pac_lag_1m, temp_mean_lag_3m", "Auto-correlation"],
    ["Rolling windows", "24", "nb_pac_rmean_3m, hdd_sum_rmean_6m", "Smoothing & volatility"],
    ["Variations", "10", "nb_pac_diff_1m, pct_change_1m", "Market acceleration"],
    ["Interactions", "4", "interact_hdd_confiance, jours_extremes", "Cross-domain hypotheses"],
    ["PAC Efficiency", "8", "cop_proxy, pac_viability_score, is_mountain", "Heat pump domain logic"],
    ["Weather", "9", "temp_mean, HDD, CDD, frost days", "Local climate drivers"],
    ["Economic", "6", "confiance_menages, IPI HVAC", "National market context"],
    ["Construction", "4", "nb_logements_autorises, surface_m2", "Real estate activity"],
    ["Socioeconomic", "8", "revenu_median, prix_m2, pct_maisons", "Department profile"],
]

add_table(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(4.5),
          len(feat_data), 4, feat_data, header_color=ACCENT_GREEN)

# PAC efficiency highlight
add_shape_rect(slide, Inches(0.4), Inches(6.0), Inches(12.5), Inches(1.2), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(0.6), Inches(6.05), Inches(12), Inches(0.4),
             "Domain Knowledge: PAC Efficiency Features", font_size=16,
             color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(0.6), Inches(6.45), Inches(12), Inches(0.7),
             "COP proxy = 4.5 - 0.08 x frost_days - 0.0005 x altitude_mean  (clipped [1.0, 5.0])\n"
             "Below -7C, air-source heat pumps lose efficiency (COP < 2.0) — "
             "mountain departments have structurally different HVAC adoption patterns.",
             font_size=13, color=DARK_GRAY)

# =========================================================================
# SLIDE 8 — ML Models & Results
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "06  Machine Learning — Models & Results", font_size=32, color=WHITE, bold=True)

# KPI cards
add_kpi_card(slide, Inches(0.8), Inches(1.4), Inches(2.5), Inches(1.3),
             "R2 = 0.989", "Ridge (Best Model)", ACCENT_GREEN)
add_kpi_card(slide, Inches(3.6), Inches(1.4), Inches(2.5), Inches(1.3),
             "RMSE = 0.93", "Test Error", MEDIUM_BLUE)
add_kpi_card(slide, Inches(6.4), Inches(1.4), Inches(2.5), Inches(1.3),
             "5 Models", "Trained & Compared", ACCENT_ORANGE)
add_kpi_card(slide, Inches(9.2), Inches(1.4), Inches(2.5), Inches(1.3),
             "551 Tests", "Passing (pytest)", RGBColor(0x7B, 0x1F, 0xA2))

# Results table
results_data = [
    ["Model", "Val RMSE", "Val R2", "Test RMSE", "Test R2", "Usage"],
    ["Ridge", "1.178", "0.9798", "0.929", "0.9885", "Baseline (best)"],
    ["LightGBM", "1.456", "0.9691", "1.283", "0.9781", "Non-linearities"],
    ["Ridge Exogenous", "1.535", "0.9657", "1.339", "0.9762", "External features"],
    ["Prophet", "—", "—", "—", "—", "Seasonality"],
    ["LSTM (PyTorch)", "—", "—", "—", "—", "DL exploration"],
]

add_table(slide, Inches(0.6), Inches(3.1), Inches(12.1), Inches(2.5),
          len(results_data), 6, results_data)

# Top features
add_text_box(slide, Inches(0.8), Inches(5.8), Inches(5), Inches(0.4),
             "Top Features (SHAP importance)", font_size=18, color=DARK_BLUE, bold=True)

top_feats = [
    "nb_installations_pac_lag_1m  — Previous month installations (auto-correlation)",
    "nb_installations_pac_diff_1m  — Month-over-month change (momentum)",
    "nb_dpe_total_rmean_3m  — 3-month rolling DPE average (trend)",
    "temp_mean_rmean_6m  — 6-month rolling temperature (seasonal climate)",
    "hdd_sum_rmean_6m  — 6-month rolling heating degree days (heating demand)",
]
add_bullet_list(slide, Inches(1.0), Inches(6.2), Inches(11), Inches(1.5),
                [f"  {f}" for f in top_feats], font_size=13, color=DARK_GRAY, spacing=Pt(4))

# =========================================================================
# SLIDE 9 — Database Schema
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "03b  Database — Star Schema", font_size=32, color=WHITE, bold=True)

# Schema boxes
def add_schema_box(slide, left, top, width, height, title, fields, color):
    shape = add_shape_rect(slide, left, top, width, height, WHITE)
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    add_shape_rect(slide, left, top, width, Inches(0.4), color)
    add_text_box(slide, left, top + Inches(0.02), width, Inches(0.36),
                 title, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.45), width - Inches(0.2),
                 height - Inches(0.5), "\n".join(fields), font_size=10, color=DARK_GRAY)

add_schema_box(slide, Inches(0.8), Inches(1.5), Inches(2.5), Inches(2.5),
               "dim_time", [
                   "date_id (PK) YYYYMM",
                   "year, month, quarter",
                   "is_heating",
                   "is_cooling"
               ], MEDIUM_BLUE)

add_schema_box(slide, Inches(0.8), Inches(4.3), Inches(2.5), Inches(2.5),
               "dim_geo", [
                   "geo_id (PK)",
                   "dept_code, dept_name",
                   "city_ref",
                   "latitude, longitude",
                   "region_code"
               ], MEDIUM_BLUE)

add_schema_box(slide, Inches(4.2), Inches(1.5), Inches(4.5), Inches(3.8),
               "fact_hvac_installations", [
                   "date_id (FK) + geo_id (FK)",
                   "nb_dpe_total (target)",
                   "nb_installations_pac (target)",
                   "nb_installations_clim",
                   "pac_per_1000_logements (normalized)",
                   "clim_per_1000_logements (normalized)",
                   "temp_mean, HDD, CDD",
                   "nb_permis_construire",
                   "UNIQUE(date_id, geo_id)"
               ], ACCENT_GREEN)

add_schema_box(slide, Inches(4.2), Inches(5.6), Inches(4.5), Inches(1.6),
               "fact_economic_context", [
                   "date_id (PK, FK) — month only",
                   "confiance_menages, climat_affaires",
                   "ipi_manufacturing, ipi_hvac_c28"
               ], ACCENT_ORANGE)

add_schema_box(slide, Inches(9.5), Inches(1.5), Inches(3.2), Inches(3.5),
               "raw_dpe", [
                   "numero_dpe (PK)",
                   "date_etablissement_dpe",
                   "code_departement",
                   "etiquette_dpe (A-G)",
                   "type_generateur_chauffage",
                   "type_generateur_froid",
                   "surface, isolation, costs",
                   "INDEX(date, dept)",
                   "~30 columns total"
               ], RGBColor(0x7B, 0x1F, 0xA2))

# Multi-engine note
add_shape_rect(slide, Inches(9.5), Inches(5.3), Inches(3.2), Inches(1.5), LIGHT_GRAY)
add_text_box(slide, Inches(9.6), Inches(5.35), Inches(3.0), Inches(1.4),
             "Multi-Engine\n\n"
             "SQLite (default)\n"
             "PostgreSQL\n"
             "SQL Server\n\n"
             "Switch via .env only",
             font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 10 — API & Dashboard
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "07  API REST & Dashboard", font_size=32, color=WHITE, bold=True)

# API section
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.4),
             "FastAPI — 6 Endpoints", font_size=20, color=MEDIUM_BLUE, bold=True)

api_data = [
    ["Endpoint", "Method", "Description"],
    ["/health", "GET", "API status, version, model info"],
    ["/predictions", "GET", "Predictions by department & horizon"],
    ["/predict", "POST", "Custom prediction with JSON params"],
    ["/data/summary", "GET", "Summary of available data"],
    ["/model/metrics", "GET", "ML metrics (RMSE, MAE, R2, MAPE)"],
    ["/departments", "GET", "List of 96 departments"],
]

add_table(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(2.8),
          len(api_data), 3, api_data, header_color=MEDIUM_BLUE)

# Dashboard section
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5), Inches(0.4),
             "Streamlit — 6 Interactive Pages", font_size=20, color=ACCENT_ORANGE, bold=True)

dashboard_data = [
    ["Page", "Description"],
    ["Home", "Overview, key metrics, data insights"],
    ["Exploration", "Stats, distributions, correlations"],
    ["Map of France", "Interactive map with metric selector"],
    ["ML Predictions", "Predictions vs actual, residuals, SHAP"],
    ["Model Comparison", "Comparison table, radar chart"],
    ["Pipeline", "Data status, pipeline launch, pCloud sync"],
]

add_table(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.8),
          len(dashboard_data), 2, dashboard_data, header_color=ACCENT_ORANGE)

# Usage examples
add_shape_rect(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.2), LIGHT_GRAY)
add_text_box(slide, Inches(0.8), Inches(5.1), Inches(5), Inches(0.4),
             "Usage Examples", font_size=18, color=DARK_BLUE, bold=True)

examples = [
    'API:  curl "http://localhost:8000/predictions?departement=69&horizon=6"',
    'Dashboard:  streamlit run app/app.py  ->  http://localhost:8501',
    'Swagger docs:  http://localhost:8000/docs  (auto-generated)',
    'Docker:  docker compose up  (API + Dashboard + optional PostgreSQL)',
]
add_bullet_list(slide, Inches(1.0), Inches(5.5), Inches(11), Inches(1.5),
                [f"  {e}" for e in examples], font_size=14, color=DARK_GRAY, spacing=Pt(4))

# =========================================================================
# SLIDE 11 — Deployment & DevOps
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "08  Deployment & DevOps", font_size=32, color=WHITE, bold=True)

# Deployment options
deploy_items = [
    ("Docker", "Multi-stage build, docker-compose (API + Dashboard + PostgreSQL)", MEDIUM_BLUE),
    ("Kubernetes", "Deployment, Service, Ingress manifests ready", ACCENT_GREEN),
    ("Airflow", "DAG for automated pipeline orchestration", ACCENT_ORANGE),
    ("Render.com", "Cloud deployment config (render.yaml)", RGBColor(0x7B, 0x1F, 0xA2)),
    ("Makefile", "make install / demo / pipeline / dashboard", DARK_GRAY),
    ("pCloud Sync", "Automated data + features upload/download", MEDIUM_BLUE),
]

for i, (title, desc, color) in enumerate(deploy_items):
    y = Inches(1.5) + i * Inches(0.75)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.8), y, Inches(2.2), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(3.2), y + Inches(0.05), Inches(9), Inches(0.45),
                 desc, font_size=14, color=DARK_GRAY)

# Quality & testing section
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(5), Inches(0.4),
             "Quality & Testing", font_size=20, color=DARK_BLUE, bold=True)

quality_items = [
    "551 unit tests (pytest) — regression + security + validation",
    "Input validation, injection prevention, dependency checks",
    "Structured logging (DEBUG / INFO / WARNING / ERROR)",
    "Centralized configuration via frozen dataclasses + .env",
]
add_bullet_list(slide, Inches(1.0), Inches(6.6), Inches(11), Inches(1.2),
                [f"  {q}" for q in quality_items], font_size=13, color=DARK_GRAY, spacing=Pt(3))

# =========================================================================
# SLIDE 12 — Certification Coverage
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "09  Data Science Lead — Certification Coverage",
             font_size=30, color=WHITE, bold=True)

cert_data = [
    ["Module", "Subject", "Implementation in Project"],
    ["M1", "Data Governance", "GDPR compliance, AI Act, data lineage, maturity assessment"],
    ["M2", "Deployment & Distributed ML", "Dockerfile, docker-compose, Kubernetes manifests, Render.com"],
    ["M3", "Database Architecture", "Star schema, OLAP design, multi-engine (SQLite/PG/MSSQL)"],
    ["M4", "Data Pipelines", "ELT pipeline, 10-step orchestrator, Airbyte alternative documented"],
    ["M5", "Automation & Workflow", "Airflow DAG, Makefile, CLI with 16 commands"],
    ["M6", "Reinforcement Learning", "Gymnasium RL demo (Q-Learning for HVAC optimization)"],
]

add_table(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(3.2),
          len(cert_data), 3, cert_data,
          header_color=RGBColor(0x7B, 0x1F, 0xA2))

# Summary statement
add_shape_rect(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.8), LIGHT_BLUE)
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.4),
             "End-to-End Coverage", font_size=20, color=DARK_BLUE, bold=True)

summary_items = [
    "Collection: 6 open data sources, 5 custom collectors with plugin architecture",
    "Processing: ELT pipeline, feature engineering (~100+ variables), outlier detection",
    "Modeling: Ridge (R2=0.989), LightGBM, Prophet, LSTM, SHAP feature importance",
    "Serving: FastAPI (6 endpoints) + Streamlit (6 pages) + Docker + Kubernetes",
    "Documentation: 3 technical docs, structured logging, 551 tests",
]
add_bullet_list(slide, Inches(1.0), Inches(5.7), Inches(11), Inches(1.5),
                [f"  {s}" for s in summary_items], font_size=14, color=DARK_GRAY, spacing=Pt(3))

# =========================================================================
# SLIDE 13 — Key Takeaways & Next Steps
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "10  Key Takeaways & Next Steps", font_size=32, color=WHITE, bold=True)

# Achievements
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.4),
             "Key Achievements", font_size=20, color=ACCENT_GREEN, bold=True)

achievements = [
    "96 departments covered with 6 open data sources",
    "~100+ engineered features with domain expertise",
    "Ridge R2=0.989 on test set (strong baseline)",
    "Full-stack delivery: API + Dashboard + Docker",
    "551 tests, professional-grade code quality",
    "100% local-first, portable, reproducible",
]
add_bullet_list(slide, Inches(1.0), Inches(1.9), Inches(5.3), Inches(3),
                [f"  {a}" for a in achievements], font_size=14, color=DARK_GRAY, spacing=Pt(5))

# Next steps
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.4),
             "Next Steps", font_size=20, color=ACCENT_ORANGE, bold=True)

next_steps = [
    "Scale DPE collection to 5-8M records (96 depts)",
    "Feature importance review & column selection",
    "Hyperparameter tuning (Ridge, LightGBM)",
    "LSTM/Transformer with sufficient data volume",
    "Candidate variables: median income, price/m2, housing stock",
    "Medium article (data storytelling)",
]
add_bullet_list(slide, Inches(7.2), Inches(1.9), Inches(5.3), Inches(3),
                [f"  {n}" for n in next_steps], font_size=14, color=DARK_GRAY, spacing=Pt(5))

# Bottom insight
add_shape_rect(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.5), LIGHT_BLUE)
add_text_box(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.4),
             "Data-Driven Insight", font_size=18, color=DARK_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9),
             "The HVAC market is driven primarily by auto-correlation (previous month installations), "
             "climate seasonality (HDD/CDD), and macro-economic confidence. Normalized targets "
             "(pac_per_1000_logements) remove departmental scale bias and enable fair comparison "
             "between urban (Rhone ~300/mo) and rural departments (~10/mo). "
             "Mountain departments show structurally different adoption patterns due to heat pump "
             "efficiency constraints below -7C.",
             font_size=14, color=DARK_GRAY)

# =========================================================================
# SLIDE 14 — Thank You
# =========================================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BLUE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
             "Thank You", font_size=52, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_shape_rect(slide, Inches(5), Inches(3.4), Inches(3), Pt(3), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "Patrice DUCLOS", font_size=24, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.4), Inches(11), Inches(0.5),
             "Senior Data Analyst  |  20 years of experience",
             font_size=16, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.4),
             "github.com/PDUCLOS/Projet-HVAC",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.7), Inches(11), Inches(0.4),
             "Data Science Lead — Jedha Bootcamp (Bac+5, RNCP Level 7)",
             font_size=14, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

# =========================================================================
# SAVE
# =========================================================================
output_path = "docs/HVAC_Market_Analysis_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
